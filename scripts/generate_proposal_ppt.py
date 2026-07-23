# -*- coding: utf-8 -*-
"""
Script tự động sinh Presentation PPTX cho Báo cáo Đề xuất Lê Thị Hồng Minh
Tuân thủ nghiêm ngặt McKinsey Presentation Playbook Standard:
- 16:9 Widescreen Layout Grid (Action Title, Breadcrumb, Core Content Body, Footnote & Source)
- High Contrast Color Palette (Slate Dark #0F172A, Card #1E293B, Text Bright White #F8FAFC)
- Action Titles theo Minto Pyramid Principle
- Embed Word-by-Slide Presenter Script vào Speaker Notes của từng slide.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------
    # MCKINSEY COLOR PALETTE (DARK EMERALD THEME)
    # -------------------------------------------------------------
    BG_DARK = RGBColor(15, 23, 42)         # Slate 900 (#0F172A)
    CARD_BG = RGBColor(30, 41, 59)         # Slate 800 (#1E293B)
    CARD_BORDER = RGBColor(51, 65, 85)     # Slate 700 (#334155)
    
    TEXT_WHITE = RGBColor(248, 250, 252)   # Slate 50 (#F8FAFC) - High Contrast Text
    TEXT_MUTED = RGBColor(148, 163, 184)   # Slate 400 (#94A3B8)
    
    ACCENT_GREEN = RGBColor(16, 185, 129)  # Emerald 500 (#10B981) - Primary Accent
    ACCENT_GOLD = RGBColor(245, 158, 11)   # Amber 500 (#F59E0B) - Secondary Accent
    ACCENT_BLUE = RGBColor(56, 189, 248)   # Sky 400 (#38BDF8) - Tertiary Accent
    
    blank_layout = prs.slide_layouts[6]
    
    def apply_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def add_slide_header(slide, action_title, breadcrumb="LEAN STARTUP PROGRAM | DOANH CHỦ LÊ THỊ HỒNG MINH"):
        tb_crumb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_c = tb_crumb.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = breadcrumb.upper()
        p_c.font.size = Pt(10)
        p_c.font.bold = True
        p_c.font.color.rgb = ACCENT_GREEN
        
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_t = tb_title.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = action_title
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

    def add_slide_footer(slide, footnote_text="Giả định kịch bản: Baseline 20 TV, Cấp GOLD, Tỷ lệ KHTT 80% - NPP 20%, F1 gánh 50%", source_text="Nguồn: Financial Simulator & Vinalink Compensation Model"):
        tb_foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4))
        tf_f = tb_foot.text_frame
        tf_f.word_wrap = True
        p_f = tf_f.paragraphs[0]
        p_f.text = f"Footnote: {footnote_text}  |  {source_text}"
        p_f.font.size = Pt(9)
        p_f.font.italic = True
        p_f.font.color.rgb = TEXT_MUTED

    def add_speaker_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # =============================================================
    # SLIDE 1: COVER SLIDE
    # =============================================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide1)
    
    tb_cover = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.0))
    tf_cov = tb_cover.text_frame
    tf_cov.word_wrap = True
    
    p0 = tf_cov.paragraphs[0]
    p0.text = "LEAN STARTUP PROGRAM 2026 — CHUYÊN ĐỀ 2 & 3"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_GREEN
    p0.space_after = Pt(16)
    
    p1 = tf_cov.add_paragraph()
    p1.text = "BÁO CÁO ĐỀ XUẤT TÁI CẤU TRÚC MỤC TIÊU &\nLỘ TRÌNH THỰC THI DÒNG TIỀN 36 THÁNG"
    p1.font.size = Pt(30)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.space_after = Pt(20)
    
    p2 = tf_cov.add_paragraph()
    p2.text = "Doanh chủ: Lê Thị Hồng Minh  |  Đơn vị tư vấn: Lean Startup Team  |  Ngày phát hành: 23/07/2026"
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MUTED

    add_speaker_notes(slide1, "Trân trọng chào chị Lê Thị Hồng Minh. Rất vui được gặp lại chị trong buổi làm việc chuyên sâu hôm nay thuộc Chuyên đề 2 & 3 của Chương trình Huấn luyện Khởi nghiệp Tinh gọn 2026. Dựa trên toàn bộ dữ liệu khảo sát hiện trạng Baseline 20 thành viên, Đội ngũ Cố vấn Lean Startup Team xin gửi tới chị Báo cáo Đề xuất Tái cấu trúc Mục tiêu và Lộ trình Vận hành 36 tháng.")

    # =============================================================
    # SLIDE 2: ACTION TITLE - ĐIỀU CHỈNH MỤC TIÊU & KỲ VỌNG
    # =============================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide2)
    add_slide_header(slide2, "1. ĐỀ XUẤT ĐIỀU CHỈNH MỤC TIÊU NĂM 3 VỀ MỐC 350M - 500M VNĐ/THÁNG ĐỂ BẢO TỒN SỨC LAO ĐỘNG VÀ TRÁNH QUÁ TẢI")
    add_slide_footer(slide2)
    
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = ACCENT_GOLD
    card1.line.width = Pt(1.5)
    
    tf1 = card1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "ĐỀ XUẤT TÁI CẤU TRÚC KỲ VỌNG THU NHẬP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.space_after = Pt(16)
    
    p1 = tf1.add_paragraph()
    p1.text = "• Mục tiêu cũ của Doanh chủ (Hồ sơ):\n  1.500.000.000 VNĐ/tháng"
    p1.font.size = Pt(13)
    p1.font.color.rgb = TEXT_MUTED
    p1.space_after = Pt(14)
    
    p2 = tf1.add_paragraph()
    p2.text = "• Đề xuất điều chỉnh mới:\n  350.000.000 — 500.000.000 VNĐ/tháng"
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN
    p2.space_after = Pt(14)
    
    p3 = tf1.add_paragraph()
    p3.text = "→ Đảm bảo tính khả thi thực tế, tạo dòng tiền thụ động bền vững mà không bị kiệt sức."
    p3.font.size = Pt(11)
    p3.font.italic = True
    p3.font.color.rgb = TEXT_WHITE

    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.0))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = ACCENT_GREEN
    card2.line.width = Pt(1.5)
    
    tf2 = card2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "LUẬN CỨ ĐỊNH LƯỢNG & NĂNG LỰC THỰC THI"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(14)
    
    p_l1 = tf2.add_paragraph()
    p_l1.text = "1. Giới hạn chi trả của Sơ đồ Nhị phân Vinalink:"
    p_l1.font.size = Pt(12)
    p_l1.font.bold = True
    p_l1.font.color.rgb = TEXT_WHITE
    
    p_l1_sub = tf2.add_paragraph()
    p_l1_sub.text = "Tại quy mô 10.000 TV, thu nhập hoạt động thuần túy tối đa ở Ambassador là ~573 triệu VNĐ/tháng. Mức 1.5 tỷ vượt quá trần chi trả của sơ đồ."
    p_l1_sub.font.size = Pt(11)
    p_l1_sub.font.color.rgb = TEXT_MUTED
    p_l1_sub.space_after = Pt(12)
    
    p_l2 = tf2.add_paragraph()
    p_l2.text = "2. Phòng ngừa rủi ro quá tải hành vi (Burnout):"
    p_l2.font.size = Pt(12)
    p_l2.font.bold = True
    p_l2.font.color.rgb = TEXT_WHITE
    
    p_l2_sub = tf2.add_paragraph()
    p_l2_sub.text = "Mục tiêu 1.5 tỷ đòi hỏi tuyển mới 140 người ngay Tháng 1, vượt quá công suất đào tạo hiện tại (Hacker score 2/5), gây đứt gãy hệ thống."
    p_l2_sub.font.size = Pt(11)
    p_l2_sub.font.color.rgb = TEXT_MUTED

    add_speaker_notes(slide2, "Xin mời chị Minh cùng nhìn vào Slide 1. Luận điểm quan trọng nhất là chúng ta cần điều chỉnh lại mục tiêu thu nhập Năm 3 từ 1.5 tỷ về 350M - 500M/tháng. Mức 1.5 tỷ vượt trần chi trả của sơ đồ nhị phân Vinalink tại quy mô 10.000 người (~573M). Mức 350M - 500M vẫn tạo ra tự do tài chính bền vững và giữ hạn mức 4 tiếng tinh gọn/ngày.")

    # =============================================================
    # SLIDE 3: BỘ CÔNG CỤ VĂN HÓA & CHẾ TÀI
    # =============================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide3)
    add_slide_header(slide3, "2. THỰC THI 3 GIÁ TRỊ CỐT LÕI VÀ BỘ QUY TẮC KỶ LUẬT ZERO TOLERANCE ĐỂ BẢO VỆ THƯƠNG HIỆU ĐỘI NGŨ")
    add_slide_footer(slide3)
    
    card_v1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    card_v1.fill.solid()
    card_v1.fill.fore_color.rgb = CARD_BG
    card_v1.line.color.rgb = ACCENT_GREEN
    card_v1.line.width = Pt(1.5)
    
    tf_v1 = card_v1.text_frame
    tf_v1.word_wrap = True
    p = tf_v1.paragraphs[0]
    p.text = "3 GIÁ TRỊ CỐT LÕI (TẢNG BĂNG VĂN HÓA)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(14)
    
    vals = [
        ("1. Đồng hành Dụng tâm (Servant Leadership):", "Lấy thành công của tuyến dưới làm thước đo. Tuyến trên hỗ trợ, không quản lý."),
        ("2. Tối giản Tinh gọn (Lean & Simplicity):", "Đóng gói bài giảng đơn giản để F5 tự sao chép được ngay."),
        ("3. Chính trực & Minh bạch (Integrity):", "Tuyệt đối không ép số ảo, không ôm hàng, không bán phá giá.")
    ]
    for t_val, d_val in vals:
        p_t = tf_v1.add_paragraph()
        p_t.text = t_val
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_d = tf_v1.add_paragraph()
        p_d.text = d_val
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_after = Pt(8)

    card_v2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.0))
    card_v2.fill.solid()
    card_v2.fill.fore_color.rgb = CARD_BG
    card_v2.line.color.rgb = ACCENT_GOLD
    card_v2.line.width = Pt(1.5)
    
    tf_v2 = card_v2.text_frame
    tf_v2.word_wrap = True
    p = tf_v2.paragraphs[0]
    p.text = "QUY TẮC KỶ LUẬT & CHẾ TÀI XỬ LÝ VI PHẠM"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.space_after = Pt(14)
    
    rules = [
        ("• Chế tài Zero Tolerance Anti-Dumping:", "Cảnh cáo lần 1; Cắt quyền bảo trợ & xử lý lần 2 đối với hành vi bán phá giá trên sàn e-commerce."),
        ("• Họp Daily Power 14:00 (15 phút Zoom):", "Đúng giờ tuyệt đối, trễ 1 phút tự giác gieo hạt 30.000 VNĐ vào quỹ chung."),
        ("• Giao tiếp Phi đồng bộ Zalo:", "Phản hồi tin nhắn công việc trong 2-4 tiếng."),
        ("• Văn hóa Không đổ lỗi (Fail Fast):", "Đối thoại bằng số liệu Simulator khi doanh số sụt giảm.")
    ]
    for r_t, r_d in rules:
        p_t = tf_v2.add_paragraph()
        p_t.text = r_t
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_d = tf_v2.add_paragraph()
        p_d.text = r_d
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_after = Pt(8)

    add_speaker_notes(slide3, "Chuyển sang Slide 2, đây chính là Bộ công cụ Tầm nhìn & Văn hóa Cốt lõi (Lean Culture Toolkit). Chúng ta thiết lập 3 Giá trị Cốt lõi: Đồng hành Dụng tâm, Tối giản Tinh gọn và Chính trực. Đồng thời áp dụng Quy tắc Zero Tolerance: Xử lý nghiêm vi phạm bán phá giá/gôm hàng để bảo vệ thương hiệu đội ngũ.")

    # =============================================================
    # SLIDE 4: BẢNG SO SÁNH 3 KỊCH BẢN
    # =============================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide4)
    add_slide_header(slide4, "3. KỊCH BẢN ỔN ĐỊNH VÀ TỐT MANG LẠI DÒNG TIỀN TÍCH LŨY 3 NĂM TỪ 5.8 TỶ ĐẾN 14.2 TỶ VNĐ")
    add_slide_footer(slide4)
    
    rows, cols = 8, 4
    t_shape = slide4.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0))
    table = t_shape.table
    table.columns[0].width = Inches(3.6)
    table.columns[1].width = Inches(2.7)
    table.columns[2].width = Inches(2.7)
    table.columns[3].width = Inches(2.7)
    
    headers = ["Chỉ số Định lượng (Tháng 36)", "Kịch bản Tệ", "Kịch bản Ổn định", "Kịch bản Tốt"]
    for col_idx, h_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = h_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN if col_idx > 0 else TEXT_WHITE
        
    data = [
        ["Quy mô tích lũy (TV)", "1.483 TV", "4.945 TV", "9.890 TV"],
        ["Doanh số tháng (CV)", "1.274.000 CV", "4.548.000 CV", "9.692.000 CV"],
        ["Danh hiệu nhóm đạt được", "Diamond", "Crown Diamond", "Ambassador"],
        ["Hoa hồng Nhóm (GVC)", "24,6 triệu VNĐ", "113,6 triệu VNĐ", "314,7 triệu VNĐ"],
        ["Matching Bonus (F1-F5)", "0 VNĐ", "10,4 triệu VNĐ", "52,0 triệu VNĐ"],
        ["TỔNG THU NHẬP THÁNG 36", "56,8 triệu VNĐ", "242,2 triệu VNĐ", "573,0 triệu VNĐ"],
        ["TỔNG TÍCH LŨY 3 NĂM", "1,03 tỷ VNĐ", "5,85 tỷ VNĐ", "14,22 tỷ VNĐ"]
    ]
    for row_idx, row_data in enumerate(data, start=1):
        is_highlight_row = row_idx in [6, 7]
        for col_idx, text_val in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text_val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG if is_highlight_row else BG_DARK
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.bold = is_highlight_row or col_idx == 0
            p.font.color.rgb = ACCENT_GOLD if (is_highlight_row and col_idx > 0) else TEXT_WHITE

    add_speaker_notes(slide4, "Xin mời chị xem Slide 3 – So sánh định lượng 3 kịch bản. Ở Kịch bản Ổn định (khuyên dùng), Tháng 36 chị đạt Crown Diamond, thu nhập 242.2M/tháng và tổng dòng tiền tích lũy 3 năm đạt 5.85 tỷ VNĐ (+ thưởng xe 1.1 tỷ). Kịch bản Tốt mang về 14.22 tỷ 3 năm.")

    # =============================================================
    # SLIDE 5: BỘ LỌC QF & MA TRẬN F1 BỔ KHUYẾT
    # =============================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide5)
    add_slide_header(slide5, "4. ÁP DỤNG BỘ LỌC 5 TIÊU CHUẨN QF ĐỂ XÂY DỰNG ĐỘI NGŨ F1 NÒNG CỐT BỔ KHUYẾT (HUSTLER - HACKER - HIPSTER)")
    add_slide_footer(slide5, footnote_text="Doanh chủ tập trung thế mạnh Quản trị (Operator 5/5), ủy quyền năng lực đào tạo và công nghệ cho F1.")
    
    roles = [
        ("F1-1: HUSTLER", "Thủ lĩnh Tuyển dụng & Kết nối", "• Chốt hợp tác nhóm Leader lớn.\n• Tạo xung lực bảo trợ từ Tháng 1-6.\n• Có uy tín cá nhân và tầm ảnh hưởng.", ACCENT_GOLD),
        ("F1-2: HACKER", "Chuyên gia Đào tạo & SOP", "• Đóng gói quy trình EMC tối giản dễ sao chép.\n• Đứng lớp huấn luyện tuyến dưới F2-F5.\n• Giúp Doanh chủ giữ hạn mức 4h/ngày.", ACCENT_GREEN),
        ("F1-3: HIPSTER", "Kỹ sư Công nghệ & Phễu", "• Thiết lập phễu tuyển dụng tự động số hóa.\n• Sản xuất video Capcut/Canva xây thương hiệu.\n• Cung cấp nguồn khách hàng tiềm năng.", ACCENT_BLUE)
    ]
    for idx, (title, subtitle, desc, color) in enumerate(roles):
        left = Inches(0.8 + idx * 4.0)
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.7), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_WHITE
        p_sub.space_after = Pt(14)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED

    add_speaker_notes(slide5, "Slide 4 thể hiện giải pháp giải phóng thời gian: Chị Minh tập trung thế mạnh Quản trị (Operator 5/5) và bổ khuyết 3 F1 (Hustler, Hacker, Hipster). Sử dụng Bộ lọc 5 tiêu chí QF (đạt 5/5 mới làm F1 nòng cốt) ngay trong OPP 2-1 để chọn đúng người.")

    # =============================================================
    # SLIDE 6: KHUYẾN NGHỊ HÀNH ĐỘNG
    # =============================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide6)
    add_slide_header(slide6, "5. KHUYẾN NGHỊ HÀNH ĐỘNG: KIÊN TRÌ KỊCH BẢN ỔN ĐỊNH, CHUẨN HÓA SOP VÀ DÙNG SIMULATOR HÀNG TUẦN")
    add_slide_footer(slide6)
    
    card_rec = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0))
    card_rec.fill.solid()
    card_rec.fill.fore_color.rgb = CARD_BG
    card_rec.line.color.rgb = ACCENT_GREEN
    card_rec.line.width = Pt(1.5)
    
    tf_rec = card_rec.text_frame
    tf_rec.word_wrap = True
    
    p = tf_rec.paragraphs[0]
    p.text = "KHUYẾN NGHỊ CHIẾN LƯỢC TỪ LEAN STARTUP TEAM"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(16)
    
    recs = [
        ("1. Kiên trì kịch bản Ổn định (Persevere):", "Điều chỉnh mục tiêu thu nhập Năm 3 về mốc 350M - 500M VNĐ/tháng để tạo nguồn thu thụ động bền vững mà không bị kiệt sức."),
        ("2. Tập trung Đóng gói SOP (Năng lực Hacker):", "Ưu tiên tuyển dụng F1 Hacker hoặc nâng cấp kỹ năng đóng gói quy trình EMC tối giản, giúp hệ thống sao chép sâu xuống tầng F5."),
        ("3. Thực thi Văn hóa Kỷ luật (Lean Culture Code):", "Duy trì nghiêm ngặt họp Daily Power 14:00 (15 phút), giao tiếp phi đồng bộ (2-4h) và cam kết Zero Tolerance với bán phá giá."),
        ("4. Quản trị bằng Số liệu Định lượng (Simulator):", "Sử dụng Bộ mô phỏng Simulator trong các buổi họp tuần để đối thoại bằng số liệu, chủ động điều hướng điểm tràn và duy trì tỷ lệ cân nhánh 70/30.")
    ]
    for title, desc in recs:
        p_t = tf_rec.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_GOLD
        
        p_d = tf_rec.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_WHITE
        p_d.space_after = Pt(10)

    add_speaker_notes(slide6, "Để kết thúc, xin mời chị nhìn vào Slide 5 – 4 Khuyến nghị hành động chiến lược: Kiên trì kịch bản Ổn định, Đóng gói bài giảng SOP, Thực thi kỷ luật văn hóa & chế tài chống phá giá, Quản trị bằng số liệu Simulator. Lean Startup Team cam kết đồng hành 2-1 cùng chị trong 30 ngày đầu.")

    out_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", "..", "04-Knowledge", "bao-cao"))
    if not os.path.exists(out_dir):
        os.makedirs(out_dir)
        
    out_path = os.path.join(out_dir, "2026-07-23-slide-de-xuat-le-thi-hong-minh.pptx")
    prs.save(out_path)
    print(f"File PPTX chuẩn McKinsey đã được sinh thành công tại: {out_path}")

if __name__ == "__main__":
    create_presentation()
